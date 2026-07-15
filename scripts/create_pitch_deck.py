"""Generate Nexus Campus Pitch Deck (11 slides) for Exam 2."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "Nexus_Campus_Pitch_Deck.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# Campus teal / charcoal (avoid purple AI cliché)
BG = RGBColor(0x0B, 0x1F, 0x2A)
ACCENT = RGBColor(0x1F, 0xA7, 0xA0)
WHITE = RGBColor(0xF5, 0xF7, 0xF8)
MUTED = RGBColor(0xB8, 0xC4, 0xCC)
CARD = RGBColor(0x12, 0x2E, 0x3C)


def set_run(run, size=20, bold=False, color=WHITE):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_bg(slide, prs):
    shape = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height  # rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def title_block(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run(run, 32, True, WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        set_run(run2, 16, False, MUTED)


def bullets(slide, items, top=1.6, left=0.55, width=12.2, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = item
        set_run(run, size, False, WHITE)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, prs)
    add_accent_bar(slide)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Portada
    s = new_slide(prs)
    box = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(3.5))
    tf = box.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "NEXUS CAMPUS"
    set_run(r, 48, True, WHITE)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Carpooling seguro para estudiantes universitarios"
    set_run(r, 22, False, ACCENT)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Startup Pitch Challenge · EPN · Julio 2026"
    set_run(r, 16, False, MUTED)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "Alisson Muñoz · Carlos Quintana"
    set_run(r, 16, False, MUTED)

    # 2 Problema
    s = new_slide(prs)
    title_block(s, "El problema", "Movilidad estudiantil cara, insegura y poco eficiente")
    bullets(
        s,
        [
            "Estudiantes de Quito pagan transporte diario alto en trayectos cortos y repetitivos.",
            "Apps generales (Uber/inDrive) no verifican pertenencia universitaria.",
            "Compartir viaje informal no tiene seguimiento, reputación ni protocolo de emergencia.",
            "Oportunidad: rutas diarias al campus ya existen; falta la capa de confianza.",
        ],
    )

    # 3 Solución
    s = new_slide(prs)
    title_block(s, "La solución", "Nexus Campus en una frase")
    bullets(
        s,
        [
            "App de viajes compartidos solo para estudiantes verificados (@epn.edu.ec).",
            "Publica ruta → solicita cupo con parada sobre la ruta → negocia precio → viaja con GPS.",
            "Seguridad: SOS, contactos de emergencia, calificaciones y aprobación de vehículo.",
            "Tarifa sugerida: $0.80 + $0.45/km (mín. $1.00), alineada al código de la app.",
        ],
    )

    # 4 Demo / Producto
    s = new_slide(prs)
    title_block(s, "Producto / Demo", "Lo que se ve en 5 minutos")
    bullets(
        s,
        [
            "Auth + onboarding institucional (pasajero / conductor).",
            "Mapa OSM + rutas OSRM + navegación en vivo con polilínea restante.",
            "Solicitudes de cupo vs notificaciones (chat, fin de viaje, SOS).",
            "Auto-notificación al llegar a parada y auto-fin al destino.",
            "Stack: Flutter · Appwrite · Riverpod · GoRouter · flutter_map.",
        ],
    )

    # 5 Modelo de negocio
    s = new_slide(prs)
    title_block(s, "Modelo de negocio", "Cómo genera dinero (proyección)")
    bullets(
        s,
        [
            "Comisión proyectada 10% sobre tarifa del viaje (ej. $0.215 en viaje promedio 3 km).",
            "Convenio institucional con universidades (piloto: $500/mes desde mes 6).",
            "Hoy el prototipo calcula y negocia precio; la comisión aún no se cobra en app.",
            "CAC estimado ~$0.34 (marketing $100/mes / 3.500 usuarios año 1).",
        ],
    )

    # 6 Mercado
    s = new_slide(prs)
    title_block(s, "Mercado", "TAM / SAM / SOM (enfoque piloto EPN)")
    bullets(
        s,
        [
            "TAM: estudiantes universitarios con smartphone en Ecuador (~1M+).",
            "SAM: universidades grandes en Quito con movilidad diaria al campus.",
            "SOM (año 1): EPN piloto → ~3.500 registrados / ~2.100 activos (mes 12).",
            "Go-to-market: campus ambassadors, redes, afiches y respaldo institucional.",
        ],
    )

    # 7 Competencia
    s = new_slide(prs)
    title_block(s, "Competencia", "Ventaja: confianza universitaria + seguridad")
    bullets(
        s,
        [
            "UTMACH Rides: similar, sin SOS ni tracking en vivo.",
            "Uber / inDrive: nacionales, sin verificación institucional.",
            "BlaBlaCar: trayectos largos, no campus.",
            "Nexus: verificación EPN + SOS + GPS + convenio institucional.",
        ],
    )

    # 8 Tracción
    s = new_slide(prs)
    title_block(s, "Tracción e hitos", "Lo construido (MVP feature-complete)")
    bullets(
        s,
        [
            "App Android operativa v1.2.13 con flujo conductor/pasajero completo.",
            "Backend Appwrite: 10 colecciones (incl. trip_locations) + storage.",
            "Repo público + documentación de negocio/técnica del examen.",
            "Siguiente: video promo, piloto cerrado EPN, medir retención semanal.",
        ],
    )

    # 9 Equipo
    s = new_slide(prs)
    title_block(s, "Equipo", "Por qué nosotros")
    bullets(
        s,
        [
            "Alisson Muñoz — producto, UX, pitch y negocio.",
            "Carlos Quintana — arquitectura Flutter/Appwrite y demo técnica.",
            "Estudiantes EPN: conocemos el dolor de movilidad del campus.",
            "Roles flexibles: CEO/Pitcher + CTO/Lead Dev (según rúbrica del curso).",
        ],
    )

    # 10 Roadmap / números
    s = new_slide(prs)
    title_block(s, "Números clave", "Proyección corregida (viaje = 1 ruta, ~3 participantes)")
    bullets(
        s,
        [
            "Tarifa promedio 3 km: $2.15 → comisión 10% = $0.215 / viaje.",
            "Costos fijos proyectados: $162/mes (marketing $100 + infra $45 + ops $17).",
            "Mes 6: ~2.880 viajes → ~$619 comisión + $500 convenio = $1.119.",
            "Punto de equilibrio ~753 viajes/mes (~94 usuarios activos).",
            "MRR mes 12 proyectado ~$1.704 | LTV 6 meses/usuario ≈ $3.4 > CAC $0.34.",
        ],
    )

    # 11 La pregunta
    s = new_slide(prs)
    title_block(s, "La pregunta", "Inversión simulada · Demo Day")
    bullets(
        s,
        [
            "Pedimos la 'inversión' completa del fondo académico (20/20).",
            "Uso de fondos (simulados): piloto EPN, marketing campus, hardening prod.",
            "Meta 90 días: 500 usuarios activos semanales y 1 convenio institucional.",
            "Nexus Campus: menos gasto, más seguridad, misma ruta al campus.",
        ],
    )

    prs.save(OUT)
    print("Saved", OUT)


if __name__ == "__main__":
    main()
